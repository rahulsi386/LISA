from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "LISA-Overview.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = "17324D"
INK = "183042"
MUTED = "607484"
WHITE = "FFFFFF"
PAPER = "F7F9FC"
BLUE = "2878B5"
BLUE_LIGHT = "DCEEFF"
TEAL = "16877A"
TEAL_LIGHT = "DDF5F2"
VIOLET = "7656A8"
VIOLET_LIGHT = "E9E2F7"
GOLD = "B7791F"
GOLD_LIGHT = "FFF0C2"
ORANGE = "C65D21"
ORANGE_LIGHT = "FFE3D3"
ROSE = "B84A70"
ROSE_LIGHT = "F7D9E3"
GREEN = "3F8F4E"
GREEN_LIGHT = "DDF3DF"
RED = "C44536"
RED_LIGHT = "FADBD8"

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=INK,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = margin
    frame.margin_right = margin
    frame.margin_top = margin
    frame.margin_bottom = margin
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, parts, x, y, w, h, size=18, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, bold, part_color in parts:
        run = paragraph.add_run()
        run.text = text
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(part_color or color)
    return box


def add_line(slide, x1, y1, x2, y2, color=MUTED, width=1.5, arrow=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_circle(slide, x, y, diameter, fill, text="", text_color=WHITE, size=18):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diameter, diameter)
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(fill)
    circle.line.color.rgb = rgb(fill)
    if text:
        add_text(
            slide,
            text,
            x,
            y,
            diameter,
            diameter,
            size=size,
            color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    return circle


def add_slide_number(slide, number):
    add_text(
        slide,
        f"{number:02d}",
        Inches(12.35),
        Inches(7.03),
        Inches(0.45),
        Inches(0.2),
        size=8,
        color="8A99A5",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_header(slide, kicker, title, subtitle=None, number=None):
    add_text(slide, kicker.upper(), Inches(0.65), Inches(0.42), Inches(5.8), Inches(0.22), size=9, color=TEAL, bold=True)
    add_text(slide, title, Inches(0.65), Inches(0.72), Inches(12.0), Inches(0.65), size=28, color=NAVY, bold=True, font=TITLE_FONT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.68), Inches(1.38), Inches(11.4), Inches(0.48), size=13, color=MUTED)
    add_rect(slide, Inches(0.65), Inches(1.86), Inches(0.72), Inches(0.05), TEAL)
    if number is not None:
        add_slide_number(slide, number)


def add_card(slide, x, y, w, h, number, title, body, accent, light, title_size=17, body_size=12):
    add_rect(slide, x, y, w, h, WHITE, "D8E1E8", radius=True)
    add_rect(slide, x, y, Inches(0.08), h, accent, accent, radius=True)
    if number:
        add_circle(slide, x + Inches(0.28), y + Inches(0.28), Inches(0.44), light, number, accent, 11)
        text_x = x + Inches(0.88)
        text_w = w - Inches(1.12)
    else:
        text_x = x + Inches(0.32)
        text_w = w - Inches(0.58)
    add_text(slide, title, text_x, y + Inches(0.27), text_w, Inches(0.34), size=title_size, color=NAVY, bold=True, font=TITLE_FONT)
    add_text(slide, body, text_x, y + Inches(0.75), text_w, h - Inches(0.93), size=body_size, color=MUTED)


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, Inches(0.0), Inches(0.0), Inches(0.14), SLIDE_H, TEAL)

    # Abstract system map: evidence enters an architect and leaves as a governed agent.
    cx, cy = Inches(10.35), Inches(3.72)
    add_circle(slide, cx - Inches(0.95), cy - Inches(0.95), Inches(1.9), TEAL, "LISA", WHITE, 24)
    orbit_nodes = [
        (Inches(8.1), Inches(1.25), BLUE, "01"),
        (Inches(11.85), Inches(1.55), VIOLET, "02"),
        (Inches(12.0), Inches(4.9), ORANGE, "03"),
        (Inches(8.05), Inches(5.4), GREEN, "04"),
    ]
    for x, y, color, label in orbit_nodes:
        add_line(slide, cx, cy, x + Inches(0.34), y + Inches(0.34), "7894A8", 1.25)
        add_circle(slide, x, y, Inches(0.68), color, label, WHITE, 10)
    add_text(slide, "REQUIRE", Inches(7.55), Inches(0.82), Inches(1.5), Inches(0.25), 8, "AFC2D2", True, align=PP_ALIGN.CENTER)
    add_text(slide, "DESIGN", Inches(11.35), Inches(1.12), Inches(1.6), Inches(0.25), 8, "AFC2D2", True, align=PP_ALIGN.CENTER)
    add_text(slide, "ASSURE", Inches(11.45), Inches(5.68), Inches(1.6), Inches(0.25), 8, "AFC2D2", True, align=PP_ALIGN.CENTER)
    add_text(slide, "DELIVER", Inches(7.55), Inches(6.15), Inches(1.3), Inches(0.25), 8, "AFC2D2", True, align=PP_ALIGN.CENTER)

    add_text(slide, "LOW CODE INTELLIGENT SYSTEM ARCHITECT", Inches(0.78), Inches(0.78), Inches(6.0), Inches(0.28), 10, "80D3CA", True)
    add_text(slide, "LISA", Inches(0.72), Inches(1.42), Inches(6.2), Inches(1.25), 54, WHITE, True, TITLE_FONT)
    add_text(slide, "From customer requirements\nto governed agent delivery", Inches(0.78), Inches(2.62), Inches(6.2), Inches(1.4), 25, WHITE, True, TITLE_FONT)
    add_text(slide, "An intelligent system architect built as automations and skills using Microsoft Scout.", Inches(0.8), Inches(4.42), Inches(5.75), Inches(0.85), 16, "D5E1EA")
    add_rect(slide, Inches(0.8), Inches(6.36), Inches(2.15), Inches(0.42), TEAL, TEAL, radius=True)
    add_text(slide, "EXECUTIVE OVERVIEW", Inches(0.8), Inches(6.36), Inches(2.15), Inches(0.42), 9, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_slide_number(slide, 1)


def what_is_lisa_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_header(slide, "What is LISA", "An architect that can carry the work", "LISA connects requirements, architecture, construction, quality, and delivery in one governed workflow.", 2)

    add_rect(slide, Inches(0.68), Inches(2.28), Inches(5.0), Inches(3.95), NAVY, NAVY, radius=True)
    add_text(slide, "LISA", Inches(1.08), Inches(2.72), Inches(4.2), Inches(0.7), 36, WHITE, True, TITLE_FONT)
    add_text(slide, "Low Code\nIntelligent\nSystem Architect", Inches(1.1), Inches(3.48), Inches(3.6), Inches(1.65), 24, "C9D8E3", True, TITLE_FONT)
    add_rect(slide, Inches(1.1), Inches(5.56), Inches(1.0), Inches(0.06), TEAL)
    add_text(slide, "Built on Microsoft Scout", Inches(1.1), Inches(5.8), Inches(3.5), Inches(0.28), 11, "80D3CA", True)

    add_card(slide, Inches(6.02), Inches(2.28), Inches(3.02), Inches(1.74), "A", "Autonomous", "Moves through analysis, design, build, evaluation, optimization, and delivery.", BLUE, BLUE_LIGHT, 16, 11)
    add_card(slide, Inches(9.34), Inches(2.28), Inches(3.02), Inches(1.74), "G", "Governed", "Pauses at human review gates and verifies consequential remote operations.", GOLD, GOLD_LIGHT, 16, 11)
    add_card(slide, Inches(6.02), Inches(4.34), Inches(3.02), Inches(1.74), "E", "Evidence-led", "Grounds decisions and tests in traceable customer evidence and current platform facts.", TEAL, TEAL_LIGHT, 16, 11)
    add_card(slide, Inches(9.34), Inches(4.34), Inches(3.02), Inches(1.74), "R", "Resumable", "Uses durable checkpoints, hashes, and manifests to recover exact workflow state.", VIOLET, VIOLET_LIGHT, 16, 11)


def why_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(slide, "Why LISA was developed", "Agent delivery has too many fragile handoffs", "LISA reduces the distance between an approved need and a demonstrable, evaluated agent.", 3)

    add_text(slide, "THE DELIVERY GAP", Inches(0.72), Inches(2.2), Inches(3.0), Inches(0.28), 10, RED, True)
    challenges = [
        ("Requirements scatter", "Critical evidence sits across documents, spreadsheets, mail, and tribal knowledge."),
        ("Architecture drifts", "Platform choices and build scope can separate from the original customer need."),
        ("Quality arrives late", "Grounding, tool use, and regressions are often tested after costly construction."),
        ("Delivery is opaque", "Without durable evidence, teams cannot explain what changed or safely resume work."),
    ]
    y = 2.62
    for index, (title, body) in enumerate(challenges, start=1):
        add_circle(slide, Inches(0.72), Inches(y), Inches(0.42), RED_LIGHT, str(index), RED, 10)
        add_text(slide, title, Inches(1.34), Inches(y - 0.01), Inches(3.4), Inches(0.3), 15, NAVY, True, TITLE_FONT)
        add_text(slide, body, Inches(1.34), Inches(y + 0.35), Inches(3.65), Inches(0.55), 10.5, MUTED)
        y += 1.08

    add_line(slide, Inches(5.26), Inches(2.28), Inches(5.26), Inches(6.68), "D5DEE5", 1)
    add_circle(slide, Inches(5.74), Inches(3.7), Inches(1.08), NAVY, "LISA", WHITE, 14)
    add_line(slide, Inches(5.25), Inches(4.24), Inches(5.72), Inches(4.24), NAVY, 2, True)
    add_line(slide, Inches(6.82), Inches(4.24), Inches(7.27), Inches(4.24), NAVY, 2, True)

    add_text(slide, "THE SYSTEM RESPONSE", Inches(7.45), Inches(2.2), Inches(3.2), Inches(0.28), 10, GREEN, True)
    responses = [
        ("Traceable interpretation", "Evidence-only requirement analysis with precise source locators."),
        ("Consistent decisions", "A repeatable platform framework and typed architecture contract."),
        ("Built-in assurance", "Evaluation gates and reversible optimization before publication."),
        ("Auditable execution", "Schema validation, hash integrity, checkpoints, and verified handoffs."),
    ]
    y = 2.62
    for index, (title, body) in enumerate(responses, start=1):
        add_circle(slide, Inches(7.45), Inches(y), Inches(0.42), GREEN_LIGHT, str(index), GREEN, 10)
        add_text(slide, title, Inches(8.07), Inches(y - 0.01), Inches(3.6), Inches(0.3), 15, NAVY, True, TITLE_FONT)
        add_text(slide, body, Inches(8.07), Inches(y + 0.35), Inches(4.3), Inches(0.55), 10.5, MUTED)
        y += 1.08


def composition_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_header(slide, "What it is made of", "One orchestrator. Nine specialist skills.", "Microsoft Scout provides the automation surface; contracts and checkpoints make the system dependable.", 4)

    add_rect(slide, Inches(0.68), Inches(2.18), Inches(12.0), Inches(0.7), NAVY, NAVY, radius=True)
    add_text(slide, "MICROSOFT SCOUT  /  AUTOMATION + SKILLS FOUNDATION", Inches(0.98), Inches(2.18), Inches(11.4), Inches(0.7), 13, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_rect(slide, Inches(0.68), Inches(3.12), Inches(2.18), Inches(2.8), TEAL, TEAL, radius=True)
    add_text(slide, "CAD", Inches(0.98), Inches(3.52), Inches(1.58), Inches(0.62), 30, WHITE, True, TITLE_FONT, PP_ALIGN.CENTER)
    add_text(slide, "ORCHESTRATOR", Inches(0.98), Inches(4.18), Inches(1.58), Inches(0.26), 9, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "Routes stages\nPresents reviews\nRecovers state", Inches(0.98), Inches(4.72), Inches(1.58), Inches(0.88), 12, "D9F4F0", False, align=PP_ALIGN.CENTER)

    groups = [
        ("UNDERSTAND", ["Requirement\nanalyzer", "Complexity\nclassifier"], BLUE, BLUE_LIGHT),
        ("CREATE", ["Solution\ndesigner", "Agent\nbuilder"], GOLD, GOLD_LIGHT),
        ("ASSURE", ["Agent\nevaluator", "Agent\noptimizer"], ROSE, ROSE_LIGHT),
        ("DELIVER", ["Artifact\ngenerator", "Artifact\npublisher", "Post-publish\ncleanup"], GREEN, GREEN_LIGHT),
    ]
    x = Inches(3.12)
    for label, skills, accent, light in groups:
        group_w = Inches(2.13) if len(skills) == 2 else Inches(2.75)
        add_text(slide, label, x, Inches(3.17), group_w, Inches(0.25), 9, accent, True, align=PP_ALIGN.CENTER)
        box_h = Inches(0.72) if len(skills) == 3 else Inches(0.92)
        gap = Inches(0.18)
        y = Inches(3.58)
        for skill in skills:
            add_rect(slide, x, y, group_w, box_h, WHITE, "D7E0E7", radius=True)
            add_rect(slide, x, y, Inches(0.08), box_h, accent, accent, radius=True)
            add_text(slide, skill, x + Inches(0.2), y, group_w - Inches(0.32), box_h, 12, NAVY, True, TITLE_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            y += box_h + gap
        x += group_w + Inches(0.28)

    add_rect(slide, Inches(3.12), Inches(6.25), Inches(9.56), Inches(0.52), "E7EDF2", "E7EDF2", radius=True)
    add_text(slide, "SHARED CONTROL PLANE  ·  CONFIG-RESOLVED PATHS  ·  SCHEMAS  ·  HASHES  ·  CHECKPOINTS", Inches(3.27), Inches(6.25), Inches(9.26), Inches(0.52), 9, MUTED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_header(slide, "How it works", "A governed path from evidence to publication", "Each stage consumes validated artifacts and produces a durable handoff for the next specialist.", 5)

    stages = [
        ("01", "Analyze", TEAL, TEAL_LIGHT),
        ("02", "Classify", VIOLET, VIOLET_LIGHT),
        ("03", "Design", BLUE, BLUE_LIGHT),
        ("04", "Build", GOLD, GOLD_LIGHT),
        ("05", "Evaluate", ORANGE, ORANGE_LIGHT),
        ("06", "Optimize", ROSE, ROSE_LIGHT),
        ("07", "Package", GREEN, GREEN_LIGHT),
        ("08", "Publish", NAVY, "DCE5EC"),
    ]
    x0 = Inches(0.68)
    y = Inches(2.55)
    stage_w = Inches(1.36)
    gap = Inches(0.14)
    centers = []
    for index, (number, title, accent, light) in enumerate(stages):
        x = x0 + index * (stage_w + gap)
        centers.append(x + stage_w // 2)
        add_rect(slide, x, y, stage_w, Inches(1.28), light, accent, radius=True)
        add_circle(slide, x + Inches(0.43), y + Inches(0.18), Inches(0.5), accent, number, WHITE, 9)
        add_text(slide, title, x + Inches(0.08), y + Inches(0.83), stage_w - Inches(0.16), Inches(0.25), 12, NAVY, True, TITLE_FONT, PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            add_line(slide, x + stage_w, y + Inches(0.64), x + stage_w + gap, y + Inches(0.64), "8A9AA7", 1.5, True)

    # Two mandatory human review gates and the evaluator/optimizer feedback loop.
    for center, label in [(centers[1], "HUMAN REVIEW"), (centers[3], "HUMAN REVIEW")]:
        diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, center - Inches(0.24), Inches(4.42), Inches(0.48), Inches(0.48))
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = rgb(GOLD_LIGHT)
        diamond.line.color.rgb = rgb(GOLD)
        add_line(slide, center, Inches(3.83), center, Inches(4.42), GOLD, 1.2)
        add_text(slide, label, center - Inches(0.58), Inches(5.03), Inches(1.16), Inches(0.25), 8, GOLD, True, align=PP_ALIGN.CENTER)

    loop_x1, loop_x2 = centers[4], centers[5]
    add_line(slide, loop_x1, Inches(2.34), loop_x1, Inches(2.05), ROSE, 1.5)
    add_line(slide, loop_x1, Inches(2.05), loop_x2, Inches(2.05), ROSE, 1.5)
    add_line(slide, loop_x2, Inches(2.05), loop_x2, Inches(2.34), ROSE, 1.5, True)
    add_text(slide, "EVIDENCE-DRIVEN RETEST LOOP", loop_x1 - Inches(0.3), Inches(1.71), Inches(2.5), Inches(0.22), 8, ROSE, True, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.68), Inches(5.8), Inches(12.0), Inches(0.66), NAVY, NAVY, radius=True)
    add_rich_text(
        slide,
        [("VALIDATE", True, "80D3CA"), (" every artifact  ·  ", False, WHITE), ("CHECKPOINT", True, "F7D692"), (" every durable boundary  ·  ", False, WHITE), ("VERIFY", True, "F4B4CA"), (" every remote write", False, WHITE)],
        Inches(1.3), Inches(6.0), Inches(10.8), Inches(0.25), 11, WHITE, PP_ALIGN.CENTER,
    )


def offers_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_header(slide, "What LISA offers", "A complete agent-delivery capability", "Not a single generator: a coordinated system for deciding, building, proving, and delivering.", 6)

    cards = [
        ("01", "Evidence intelligence", "Multi-format extraction\nTraceable findings and gaps\nDeterministic customer analysis", TEAL, TEAL_LIGHT),
        ("02", "Architecture decisions", "Platform and harness selection\nComplexity and coverage scoring\nBuildable topology and diagrams", VIOLET, VIOLET_LIGHT),
        ("03", "Agent construction", "Harness-aware implementation\nLive-state reconciliation\nInstructions, package, and evidence", GOLD, GOLD_LIGHT),
        ("04", "Quality assurance", "Grounded evaluation dataset\nFour deployment gates\nReversible evidence-led optimization", ORANGE, ORANGE_LIGHT),
        ("05", "Customer artifacts", "Solution document\nInteractive execution tree\nPresentation-ready architecture", BLUE, BLUE_LIGHT),
        ("06", "Controlled delivery", "Idempotent SharePoint publication\nFresh remote verification\nFingerprint-bound local cleanup", GREEN, GREEN_LIGHT),
    ]
    positions = [
        (0.68, 2.22), (4.48, 2.22), (8.28, 2.22),
        (0.68, 4.45), (4.48, 4.45), (8.28, 4.45),
    ]
    for card, (x, y) in zip(cards, positions):
        number, title, body, accent, light = card
        add_card(slide, Inches(x), Inches(y), Inches(3.48), Inches(1.88), number, title, body, accent, light, 16, 11)


def controls_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_text(slide, "TRUST & CONTROL", Inches(0.7), Inches(0.5), Inches(3.5), Inches(0.24), 9, "80D3CA", True)
    add_text(slide, "Autonomous does not mean unsupervised", Inches(0.7), Inches(0.91), Inches(11.7), Inches(0.68), 29, WHITE, True, TITLE_FONT)
    add_text(slide, "LISA is designed to move independently inside explicit evidence, identity, and approval boundaries.", Inches(0.72), Inches(1.66), Inches(10.8), Inches(0.42), 14, "C8D6E0")

    controls = [
        ("01", "Human review gates", "Classification and build require Accept, Revise, or Cancel decisions.", GOLD),
        ("02", "Environment identity", "Tenant, environment, agent, and browser identity are verified before writes.", TEAL),
        ("03", "Artifact integrity", "Schemas, hashes, manifests, and canonical paths protect every handoff.", BLUE),
        ("04", "Exact recovery", "Checkpoints resume the recorded phase instead of guessing from folders.", VIOLET),
        ("05", "Independent evaluation", "Evaluator owns deployment gates; optimizer cannot score its own changes.", ROSE),
        ("06", "Consent-bound cleanup", "Local deletion requires fresh inventory and the exact confirmation phrase.", GREEN),
    ]
    positions = [(0.72, 2.5), (4.48, 2.5), (8.24, 2.5), (0.72, 4.55), (4.48, 4.55), (8.24, 4.55)]
    for (number, title, body, accent), (x, y) in zip(controls, positions):
        add_rect(slide, Inches(x), Inches(y), Inches(3.42), Inches(1.62), "223F5A", "36556F", radius=True)
        add_circle(slide, Inches(x + 0.27), Inches(y + 0.27), Inches(0.42), accent, number, WHITE, 9)
        add_text(slide, title, Inches(x + 0.87), Inches(y + 0.27), Inches(2.25), Inches(0.3), 14, WHITE, True, TITLE_FONT)
        add_text(slide, body, Inches(x + 0.87), Inches(y + 0.72), Inches(2.25), Inches(0.62), 10.5, "C8D6E0")
    add_slide_number(slide, 7)


def close_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, Inches(0.0), Inches(0.0), Inches(0.14), SLIDE_H, TEAL)
    add_text(slide, "LISA", Inches(0.82), Inches(0.68), Inches(3.2), Inches(0.55), 31, NAVY, True, TITLE_FONT)
    add_text(slide, "A low-code intelligent system architect\nfor governed agent delivery", Inches(0.82), Inches(1.62), Inches(7.4), Inches(1.3), 29, NAVY, True, TITLE_FONT)
    add_text(slide, "Requirements become evidence. Evidence becomes architecture.\nArchitecture becomes an evaluated, deliverable agent.", Inches(0.84), Inches(3.25), Inches(6.7), Inches(0.9), 16, MUTED)

    steps = [
        ("1", "Provide", "requirements + evaluation material", TEAL),
        ("2", "Configure", "target environment + registry", BLUE),
        ("3", "Review", "classification + constructed agent", GOLD),
        ("4", "Receive", "verified agent + customer artifacts", GREEN),
    ]
    x = Inches(0.84)
    for number, verb, detail, accent in steps:
        add_circle(slide, x, Inches(5.25), Inches(0.48), accent, number, WHITE, 10)
        add_text(slide, verb, x + Inches(0.64), Inches(5.22), Inches(1.25), Inches(0.28), 13, NAVY, True, TITLE_FONT)
        add_text(slide, detail, x + Inches(0.64), Inches(5.57), Inches(1.78), Inches(0.52), 9.5, MUTED)
        x += Inches(3.02)

    add_rect(slide, Inches(8.65), Inches(0.92), Inches(3.74), Inches(3.75), NAVY, NAVY, radius=True)
    add_circle(slide, Inches(9.65), Inches(1.57), Inches(1.74), TEAL, "LISA", WHITE, 23)
    add_text(slide, "REQUIRE", Inches(8.96), Inches(3.75), Inches(0.82), Inches(0.22), 8, "80D3CA", True, align=PP_ALIGN.CENTER)
    add_text(slide, "DESIGN", Inches(9.83), Inches(3.75), Inches(0.82), Inches(0.22), 8, "A8CDEA", True, align=PP_ALIGN.CENTER)
    add_text(slide, "ASSURE", Inches(10.7), Inches(3.75), Inches(0.82), Inches(0.22), 8, "F4B4CA", True, align=PP_ALIGN.CENTER)
    add_text(slide, "DELIVER", Inches(11.57), Inches(3.75), Inches(0.82), Inches(0.22), 8, "B8DFBF", True, align=PP_ALIGN.CENTER)
    add_text(slide, "Microsoft Scout automation + skills", Inches(8.96), Inches(4.18), Inches(3.12), Inches(0.27), 10, "C8D6E0", True, align=PP_ALIGN.CENTER)
    add_slide_number(slide, 8)


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "LISA - Low Code Intelligent System Architect"
    prs.core_properties.subject = "Executive overview of LISA"
    prs.core_properties.author = "LISA"
    prs.core_properties.keywords = "LISA, Microsoft Scout, Copilot, agent delivery"

    cover_slide(prs)
    what_is_lisa_slide(prs)
    why_slide(prs)
    composition_slide(prs)
    pipeline_slide(prs)
    offers_slide(prs)
    controls_slide(prs)
    close_slide(prs)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_deck())