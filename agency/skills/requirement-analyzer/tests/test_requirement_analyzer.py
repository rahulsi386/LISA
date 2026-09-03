from __future__ import annotations

import importlib.util
import io
import json
import shutil
import subprocess
import sys
import tempfile
import unittest
import warnings
import zipfile
from email.message import EmailMessage
from pathlib import Path


SKILL_ROOT = Path(__file__).resolve().parents[1]
SCRIPT = SKILL_ROOT / "scripts" / "requirement_analyzer.py"
FIXTURE = SKILL_ROOT / "tests" / "fixtures" / "basic"

spec = importlib.util.spec_from_file_location("requirement_analyzer", SCRIPT)
assert spec and spec.loader
analyzer = importlib.util.module_from_spec(spec)
spec.loader.exec_module(analyzer)


class RequirementAnalyzerTests(unittest.TestCase):
    maxDiff = None

    def run_cli(self, *arguments: str, expected: int = 0) -> dict:
        completed = subprocess.run(
            [sys.executable, str(SCRIPT), *arguments],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertEqual(
            expected,
            completed.returncode,
            msg=f"stdout:\n{completed.stdout}\nstderr:\n{completed.stderr}",
        )
        payload = completed.stdout.strip() or completed.stderr.strip()
        return json.loads(payload)

    def prepare_fixture(self, temporary: Path) -> tuple[dict, dict]:
        source_base = temporary / "source"
        shutil.copytree(FIXTURE, source_base)
        legacy_requirements = source_base / "Requirements"
        requirements = source_base / "requirements"
        legacy_requirements.rename(requirements)
        config = source_base / "lisa-config.json"
        config.write_text(json.dumps({"basePath": "."}), encoding="utf-8")

        message = EmailMessage()
        message["Subject"] = "Requirement example"
        message.set_content(
            "The assistant should preserve source traceability. "
            "Policy reference: https://example.com/policy"
        )
        message.add_attachment(
            b"Embedded requirement evidence.",
            maintype="text",
            subtype="plain",
            filename="attachment.txt",
        )
        (requirements / "example.eml").write_bytes(message.as_bytes())

        output = source_base / "output"
        prepared = self.run_cli(
            "prepare",
            "--config",
            str(config),
            "--workers",
            "2",
        )
        run = json.loads(Path(prepared["run"]).read_text(encoding="utf-8"))
        manifest = json.loads(Path(prepared["manifest"]).read_text(encoding="utf-8"))
        self.assertTrue(Path(prepared["review_pack"]).exists())
        analysis = output / "analysis"
        self.assertEqual(output.resolve(), Path(run["temp_output_path"]))
        self.assertEqual(analysis.resolve(), Path(run["output_root"]))
        self.assertEqual(["analysis"], sorted(item.name for item in output.iterdir()))
        self.assertTrue(Path(prepared["run"]).is_relative_to(analysis))
        self.assertEqual(analysis.resolve(), Path(prepared["target_markdown"]).parent)
        return run, manifest

    @staticmethod
    def make_draft_ledger(run: dict, manifest: dict) -> dict:
        findings = []
        source_annotations = []
        next_number = 1

        def add_finding(
            kind: str,
            statement: str,
            status: str,
            confidence: str,
            source_id: str,
            locator: str,
            evidence_type: str,
        ) -> str:
            nonlocal next_number
            identifier = f"D-{next_number:03d}"
            next_number += 1
            findings.append(
                {
                    "finding_id": identifier,
                    "kind": kind,
                    "statement": statement,
                    "status": status,
                    "confidence": confidence,
                    "evidence": [
                        {
                            "source_id": source_id,
                            "locator": locator,
                            "evidence_type": evidence_type,
                        }
                    ],
                }
            )
            return identifier

        for source in manifest["sources"]:
            identifier = add_finding(
                "Observed fact",
                f"The source inventory contains {source['relative_path']}.",
                "Current",
                "Confirmed",
                source["source_id"],
                "filesystem metadata",
                "metadata",
            )
            source_annotations.append(
                {
                    "source_id": source["source_id"],
                    "role": "Requirement evidence",
                    "classification": "Folder file",
                    "finding_ids": [identifier],
                }
            )

        first_source = manifest["sources"][0]["source_id"]
        conversational = add_finding(
            "Derived classification",
            "The requested assistant interaction is conversational.",
            "Confirmed",
            "Confirmed",
            first_source,
            "line 1",
            "derived",
        )
        platform_gap = add_finding(
            "Analyst-identified gap",
            "No permitted agent-development platform is evidenced.",
            "Not evidenced",
            "Missing",
            "CORPUS",
            "platform completeness check across the extracted corpus",
            "absence_check",
        )
        knowledge_gap = add_finding(
            "Analyst-identified gap",
            "No source qualifies as a runtime knowledge source.",
            "Not evidenced",
            "Missing",
            "CORPUS",
            "knowledge-source qualification check across the extracted corpus",
            "absence_check",
        )
        integration_gap = add_finding(
            "Analyst-identified gap",
            "No integration is evidenced.",
            "Not evidenced",
            "Missing",
            "CORPUS",
            "integration completeness check across the extracted corpus",
            "absence_check",
        )
        autonomous_gap = add_finding(
            "Analyst-identified gap",
            "Autonomous behavior is not evidenced.",
            "Not evidenced",
            "Missing",
            "CORPUS",
            "agentic-behavior completeness check across the extracted corpus",
            "absence_check",
        )
        multi_agent_gap = add_finding(
            "Analyst-identified gap",
            "Child-agent or multi-agent behavior is not evidenced.",
            "Not evidenced",
            "Missing",
            "CORPUS",
            "agentic-behavior completeness check across the extracted corpus",
            "absence_check",
        )
        oversight_gap = add_finding(
            "Analyst-identified gap",
            "Human handoff or oversight is not evidenced.",
            "Not evidenced",
            "Missing",
            "CORPUS",
            "agentic-behavior completeness check across the extracted corpus",
            "absence_check",
        )

        first_inventory_finding = source_annotations[0]["finding_ids"][0]
        section_names = [
            "Executive Summary",
            "Problem Statement",
            "Current State",
            "Desired Future State",
            "Goals",
            "Success Criteria",
            "Metrics and Baselines",
            "Data Sources",
            "Data Types",
            "Dependencies and Constraints",
            "Solution Components",
            "Scope and Delivery Phases",
            "Gaps and Conflicts",
        ]
        sections = {
            name: [
                {
                    "type": "paragraph",
                    "text": f"{name} is grounded in the inventoried requirement evidence.",
                    "finding_ids": [first_inventory_finding],
                }
            ]
            for name in section_names
        }

        return {
            "schema_version": "1.0",
            "run_id": run["run_id"],
            "findings": findings,
            "source_annotations": source_annotations,
            "referred_artifacts": [],
            "manual_reviews": [],
            "knowledge_sources": [],
            "knowledge_source_notes": [
                {
                    "text": (
                        "No confirmed or strong-candidate runtime knowledge sources "
                        "were evidenced."
                    ),
                    "finding_ids": [knowledge_gap],
                }
            ],
            "platforms": [],
            "platform_absence_finding_ids": [platform_gap],
            "integrations": [],
            "integration_notes": [
                {
                    "text": "No integrations were evidenced.",
                    "finding_ids": [integration_gap],
                }
            ],
            "agentic_behaviors": [
                {
                    "behavior": "Conversational",
                    "requirement_status": "Confirmed",
                    "evidenced_behavior": "Users ask the assistant questions.",
                    "trigger_decision_handoff": "User prompt",
                    "gaps": "Channel and response targets are unspecified.",
                    "finding_ids": [conversational],
                },
                {
                    "behavior": "Autonomous",
                    "requirement_status": "Not evidenced",
                    "evidenced_behavior": "No autonomous behavior is described.",
                    "trigger_decision_handoff": "Not evidenced",
                    "gaps": "No autonomous requirements are supplied.",
                    "finding_ids": [autonomous_gap],
                },
                {
                    "behavior": "Child-Agent/Multi-Agent",
                    "requirement_status": "Not evidenced",
                    "evidenced_behavior": "No delegation or multi-agent behavior is described.",
                    "trigger_decision_handoff": "Not evidenced",
                    "gaps": "No multi-agent requirements are supplied.",
                    "finding_ids": [multi_agent_gap],
                },
                {
                    "behavior": "Human Handoff/oversight",
                    "requirement_status": "Not evidenced",
                    "evidenced_behavior": "No human handoff is described.",
                    "trigger_decision_handoff": "Not evidenced",
                    "gaps": "No oversight requirement is supplied.",
                    "finding_ids": [oversight_gap],
                },
            ],
            "sections": sections,
        }

    def normalize_draft(
        self, run: dict, draft: dict, name: str = "test"
    ) -> Path:
        draft_path = Path(run["run_directory"]) / f"ledger.{name}.draft.json"
        draft_path.write_text(json.dumps(draft, indent=2), encoding="utf-8")
        normalized_path = (
            Path(run["run_directory"]) / f"ledger.{name}.normalized.json"
        )
        self.run_cli(
            "normalize",
            "--run",
            str(Path(run["run_directory"]) / "run.json"),
            "--ledger",
            str(draft_path),
            "--output",
            str(normalized_path),
        )
        return normalized_path

    def test_resolves_only_exact_configured_root(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            base = Path(directory)
            exact = base / "requirements"
            exact.mkdir()
            self.assertEqual(exact.resolve(), analyzer.resolve_requirements_root(exact))

        with tempfile.TemporaryDirectory() as directory:
            base = Path(directory)
            with self.assertRaises(analyzer.AnalyzerError):
                analyzer.resolve_requirements_root(base)

    def test_native_office_and_email_extraction(self) -> None:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation

        docx = io.BytesIO()
        document = Document()
        document.add_paragraph("Document requirement")
        document.save(docx)
        result = analyzer._extract_bytes(
            "requirement.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            docx.getvalue(),
        )
        self.assertEqual("manual-review-required", result["status"])
        self.assertEqual("render:all-pages", result["review_targets"][0]["target_key"])
        self.assertIn("Document requirement", analyzer._collect_extracted_text(result))

        xlsx = io.BytesIO()
        workbook = Workbook()
        workbook.active["A1"] = "Spreadsheet requirement"
        workbook.save(xlsx)
        result = analyzer._extract_bytes(
            "requirement.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            xlsx.getvalue(),
        )
        self.assertEqual("complete", result["status"])
        self.assertEqual("xlsx-stream-profile", result["method"])
        self.assertEqual(0, len(result["review_targets"]))
        self.assertEqual(1, result["metadata"]["rows_scanned"])
        self.assertIn("Spreadsheet requirement", analyzer._collect_extracted_text(result))
        row_number, values = analyzer._parse_xlsx_sample_row(
            b'<row r="1" x14ac:dyDescent="0.25"><c r="A1" t="inlineStr">'
            b"<is><t>Header</t></is></c></row>"
        )
        self.assertEqual("1", row_number)
        self.assertEqual("Header", values[0]["value"])

        pptx = io.BytesIO()
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Presentation requirement"
        presentation.save(pptx)
        result = analyzer._extract_bytes(
            "requirement.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            pptx.getvalue(),
        )
        self.assertEqual("manual-review-required", result["status"])
        self.assertEqual(1, len(result["review_targets"]))
        self.assertIn("Presentation requirement", analyzer._collect_extracted_text(result))

    def test_visual_and_unsafe_archive_require_review(self) -> None:
        from PIL import Image

        image_bytes = io.BytesIO()
        Image.new("RGB", (20, 20), "white").save(image_bytes, format="PNG")
        image = analyzer._extract_bytes("diagram.png", "image/png", image_bytes.getvalue())
        self.assertEqual("manual-review-required", image["status"])
        self.assertEqual(1, len(image["review_targets"]))

        archive_bytes = io.BytesIO()
        with zipfile.ZipFile(archive_bytes, "w") as archive:
            archive.writestr("../escape.txt", "unsafe")
        archive = analyzer._extract_bytes(
            "unsafe.zip", "application/zip", archive_bytes.getvalue()
        )
        self.assertEqual("failed", archive["status"])
        self.assertIn("Unsafe archive path", archive["warnings"][0])

        drive_archive = io.BytesIO()
        with zipfile.ZipFile(drive_archive, "w") as archive_file:
            archive_file.writestr("C:/escape.txt", "unsafe")
        archive = analyzer._extract_bytes(
            "drive.zip", "application/zip", drive_archive.getvalue()
        )
        self.assertEqual("failed", archive["status"])
        self.assertIn("Unsafe archive path", archive["warnings"][0])

        duplicate_archive = io.BytesIO()
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", UserWarning)
            with zipfile.ZipFile(duplicate_archive, "w") as archive_file:
                archive_file.writestr("duplicate.txt", "first")
                archive_file.writestr("duplicate.txt", "second")
        archive = analyzer._extract_bytes(
            "duplicate.zip", "application/zip", duplicate_archive.getvalue()
        )
        self.assertEqual("failed", archive["status"])
        self.assertIn("Duplicate archive path", archive["warnings"][0])

        from pypdf import PdfWriter

        pdf_bytes = io.BytesIO()
        writer = PdfWriter()
        writer.add_blank_page(width=100, height=100)
        writer.write(pdf_bytes)
        pdf = analyzer._extract_bytes(
            "requirement.pdf", "application/pdf", pdf_bytes.getvalue()
        )
        self.assertEqual("manual-review-required", pdf["status"])
        self.assertEqual("render:page:1", pdf["review_targets"][0]["target_key"])

    def test_prepare_uses_cache_and_records_embedded_content(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            self.assertEqual(2, manifest["source_count"])
            self.assertEqual(0, run["review_target_count"])
            self.assertTrue(
                all(item["extraction_status"] == "complete" for item in manifest["sources"])
            )

            second = self.run_cli(
                "prepare",
                "--config",
                str(temporary / "source" / "lisa-config.json"),
            )
            self.assertEqual(2, second["cache_hits"])

    def test_normalize_render_and_validate_end_to_end(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            draft = self.make_draft_ledger(run, manifest)
            normalized_path = self.normalize_draft(run, draft)
            normalized = json.loads(normalized_path.read_text(encoding="utf-8"))
            self.assertTrue(
                all(
                    analyzer.FINAL_FINDING_PATTERN.fullmatch(item["finding_id"])
                    for item in normalized["findings"]
                )
            )

            rendered = self.run_cli(
                "render",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(normalized_path),
            )
            markdown = Path(rendered["markdown"])
            self.assertTrue(markdown.exists())
            self.assertTrue(Path(rendered["ledger"]).exists())
            self.assertTrue(Path(rendered["manifest"]).exists())
            rendered_run = json.loads(
                (Path(run["run_directory"]) / "run.json").read_text(encoding="utf-8")
            )
            self.assertEqual("rendered_pending_validation", rendered_run["status"])

            validated = self.run_cli(
                "validate",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(normalized_path),
                "--markdown",
                str(markdown),
            )
            self.assertEqual("passed", validated["status"])
            validated_run = json.loads(
                (Path(run["run_directory"]) / "run.json").read_text(encoding="utf-8")
            )
            self.assertEqual("validated", validated_run["status"])

            cached = self.run_cli(
                "prepare",
                "--config",
                str(temporary / "source" / "lisa-config.json"),
            )
            self.assertTrue(cached["analysis_cache_hit"])
            self.assertTrue(Path(cached["reused_ledger"]).exists())
            fast = self.run_cli(
                "publish",
                "--run",
                cached["run"],
                "--ledger",
                cached["reused_ledger"],
            )
            self.assertEqual("validated", fast["status"])
            self.assertTrue(fast["analysis_cache_hit"])
            self.assertLess(fast["duration_seconds"], 30)

    def test_render_rejects_source_changes_after_prepare(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            draft = self.make_draft_ledger(run, manifest)
            normalized_path = self.normalize_draft(run, draft, "stale-source")
            (temporary / "source" / "requirements" / "brief.txt").write_text(
                "Changed requirement evidence.", encoding="utf-8"
            )
            result = self.run_cli(
                "render",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(normalized_path),
                expected=2,
            )
            self.assertIn("changed after preparation", result["error"])

    def test_validation_rejects_tampered_published_artifacts(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            draft = self.make_draft_ledger(run, manifest)
            normalized_path = self.normalize_draft(run, draft, "tamper")
            rendered = self.run_cli(
                "render",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(normalized_path),
            )
            markdown = Path(rendered["markdown"])
            markdown.write_text(
                markdown.read_text(encoding="utf-8") + "\nTampered prose.\n",
                encoding="utf-8",
            )
            result = self.run_cli(
                "validate",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(normalized_path),
                "--markdown",
                str(markdown),
                expected=2,
            )
            self.assertIn(
                "Published Markdown differs from deterministic ledger rendering",
                result["errors"],
            )

    def test_knowledge_source_provenance_is_source_specific(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            draft = self.make_draft_ledger(run, manifest)
            source_by_name = {
                item["relative_path"]: item["source_id"] for item in manifest["sources"]
            }
            email_source = source_by_name["example.eml"]
            brief_source = source_by_name["brief.txt"]
            draft_id = f"D-{len(draft['findings']) + 1:03d}"
            draft["findings"].append(
                {
                    "finding_id": draft_id,
                    "kind": "Observed fact",
                    "statement": "The email references an external policy URL.",
                    "status": "Current",
                    "confidence": "Confirmed",
                    "evidence": [
                        {
                            "source_id": email_source,
                            "locator": "body part 1 (text/plain; utf-8)",
                            "evidence_type": "explicit",
                            "quote": "https://example.com/policy",
                        }
                    ],
                }
            )
            draft["knowledge_sources"].append(
                {
                    "name": "Policy site",
                    "classification": "Strong Candidate Knowledge Source",
                    "location": "https://example.com/policy",
                    "source_id": brief_source,
                    "hosting_type": "Public website",
                    "content_type": "Policy",
                    "structure": "Web page",
                    "grounding_purpose": "Answer policy questions",
                    "authority_priority": "Unspecified",
                    "access_behavior": "Unspecified",
                    "ownership_freshness": "Unspecified",
                    "intended_use": "Runtime grounding",
                    "finding_ids": [draft_id],
                }
            )
            normalized_path = self.normalize_draft(run, draft, "wrong-provenance")
            result = self.run_cli(
                "render",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(normalized_path),
                expected=2,
            )
            self.assertIn("different source", result["error"])

    def test_schema_rejects_unknown_finding_status(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            draft = self.make_draft_ledger(run, manifest)
            draft["findings"][0]["status"] = "Unknown status"
            draft_path = Path(run["run_directory"]) / "ledger.invalid.draft.json"
            draft_path.write_text(json.dumps(draft, indent=2), encoding="utf-8")
            result = self.run_cli(
                "normalize",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(draft_path),
                expected=2,
            )
            self.assertIn("Evidence ledger schema error", result["error"])

    def test_schema_requires_observed_visual_evidence_details(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            draft = self.make_draft_ledger(run, manifest)
            draft["manual_reviews"] = [
                {
                    "target_id": "REV-EXAMPLE",
                    "status": "complete",
                    "method": "visual inspection",
                    "coverage": "entire image",
                    "notes": "Reviewed.",
                    "result": "evidence-observed",
                    "observations": [],
                }
            ]
            draft_path = Path(run["run_directory"]) / "ledger.review.draft.json"
            draft_path.write_text(json.dumps(draft, indent=2), encoding="utf-8")
            result = self.run_cli(
                "normalize",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(draft_path),
                expected=2,
            )
            self.assertIn("Evidence ledger schema error", result["error"])

    def test_absence_claims_require_corpus_gap_findings(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            temporary = Path(directory)
            run, manifest = self.prepare_fixture(temporary)
            draft = self.make_draft_ledger(run, manifest)
            draft["platform_absence_finding_ids"] = [
                draft["source_annotations"][0]["finding_ids"][0]
            ]
            normalized_path = self.normalize_draft(run, draft, "invalid-absence")
            result = self.run_cli(
                "render",
                "--run",
                str(Path(run["run_directory"]) / "run.json"),
                "--ledger",
                str(normalized_path),
                expected=2,
            )
            self.assertIn("Platform absence must cite", result["error"])

    def test_configured_knowledge_sources_are_prepopulated(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            base = Path(directory)
            requirements = base / "requirements"
            requirements.mkdir()
            (requirements / "brief.txt").write_text(
                "Users require a searchable assistant.", encoding="utf-8"
            )
            config = base / "lisa-config.json"
            config.write_text(
                json.dumps(
                    {
                        "basePath": ".",
                        "knowledgeSources": [
                            {
                                "name": "Procurement Policy",
                                "path": "https://contoso.sharepoint.com/policy",
                            }
                        ],
                    }
                ),
                encoding="utf-8",
            )
            prepared = self.run_cli("prepare", "--config", str(config))
            draft = json.loads(
                Path(prepared["ledger_draft"]).read_text(encoding="utf-8")
            )
            self.assertEqual(
                "Configured Knowledge Source",
                draft["knowledge_sources"][0]["classification"],
            )
            self.assertEqual("CONFIG", draft["knowledge_sources"][0]["source_id"])
            self.assertEqual(
                "https://contoso.sharepoint.com/policy",
                draft["knowledge_sources"][0]["location"],
            )

    def test_absolute_base_path_is_supported(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            base = Path(directory)
            requirements = base / "requirements"
            requirements.mkdir()
            (requirements / "brief.txt").write_text("Requirement", encoding="utf-8")
            config = base / "lisa-config.json"
            config.write_text(
                json.dumps({"basePath": str(base.resolve())}),
                encoding="utf-8",
            )
            result = self.run_cli(
                "prepare",
                "--config",
                str(config),
            )
            self.assertEqual("prepared", result["status"])
            self.assertTrue((base / "output" / "analysis").is_dir())

    def test_default_temp_output_uses_analysis_child(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            base = Path(directory)
            requirements = base / "requirements"
            requirements.mkdir()
            (requirements / "brief.txt").write_text(
                "Users require a searchable assistant.", encoding="utf-8"
            )
            config = base / "lisa-config.json"
            config.write_text(json.dumps({"basePath": "."}), encoding="utf-8")
            prepared = self.run_cli(
                "prepare",
                "--config",
                str(config),
            )
            run = json.loads(Path(prepared["run"]).read_text(encoding="utf-8"))
            self.assertEqual((base / "output").resolve(), Path(run["temp_output_path"]))
            self.assertEqual(
                (base / "output" / "analysis").resolve(),
                Path(run["output_root"]),
            )
            self.assertEqual(
                ["analysis"],
                sorted(item.name for item in (base / "output").iterdir()),
            )


if __name__ == "__main__":
    unittest.main()
